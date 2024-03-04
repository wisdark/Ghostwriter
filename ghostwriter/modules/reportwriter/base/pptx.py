
import io
import logging

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.parts.presentation import PresentationPart
from pptx.exc import PackageNotFoundError
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.enum.text import MSO_AUTO_SIZE

from ghostwriter.commandcenter.models import CompanyInformation
from ghostwriter.modules.reportwriter.base.base import ExportBase
from ghostwriter.modules.reportwriter.richtext.pptx import HtmlToPptxWithEvidence

logger = logging.getLogger(__name__)


class ExportBasePptx(ExportBase):
    """
    Base class for exporting Pptx (PowerPoint) files

    Subclasses should override `run` to add slides to the `ppt_presentation` field, using `process_rich_text_pptx`
    to template and convert rich text fields, then return `super().run()` to save and return the presentation.
    """
    ppt_presentation: PresentationPart
    company_config: CompanyInformation

    def __init__(self, object, template_loc=None):
        super().__init__(object)

        try:
            self.ppt_presentation = Presentation(template_loc)
        except ValueError:
            logger.exception(
                "Failed to load the provided template document because it is not a PowerPoint file: %s",
                template_loc,
            )
            raise
        except PackageNotFoundError:
            logger.exception(
                "Failed to load the provided template document because file could not be found: %s",
                template_loc,
            )
            raise
        except Exception:
            logger.exception(
                "Failed to load the provided template document for unknown reason: %s",
                template_loc,
            )
            raise

        self.company_config = CompanyInformation.get_solo()

    def process_rich_text_pptx(self, text, slide, shape, template_vars, evidences):
        """
        Converts HTML from the TinyMCE rich text editor and inserts it into the passed in slide and shape
        """
        text = self.preprocess_rich_text(text, template_vars)
        try:
            HtmlToPptxWithEvidence.run(
                text,
                slide=slide,
                shape=shape,
                evidences=evidences,
            )
        except:
            # Log input text to help diagnose errors
            logger.warning("Input text: %r", text)
            raise

    def run(self):
        out = io.BytesIO()
        self.ppt_presentation.save(out)
        return out


def add_slide_number(txtbox):
    """
    Add a slide number to the provided textbox. Ideally, the textbox should be the slide layout's slide
    number placeholder to match the template.

    Ref: https://stackoverflow.com/a/55816723
    """
    # Get a textbox's paragraph element
    par = txtbox.text_frame.paragraphs[0]._p

    # The slide number is actually a field, so we add a `fld` element to the paragraph
    # The number enclosed in the `a:t` element is the slide number and should auto-update on load/shuffle
    fld_xml = (
        '<a:fld %s id="{1F4E2DE4-8ADA-4D4E-9951-90A1D26586E7}" type="slidenum">\n'
        '  <a:rPr lang="en-US" smtClean="0"/>\n'
        "  <a:t>2</a:t>\n"
        "</a:fld>\n" % nsdecls("a")
    )
    fld = parse_xml(fld_xml)
    par.append(fld)


def clone_placeholder(slide, slide_layout, placeholder_idx):
    """
    Clone a placeholder from the slide master and return the layout and the new shape.
    """
    layout_placeholder = slide_layout.placeholders[placeholder_idx]
    slide.shapes.clone_placeholder(layout_placeholder)

    # The cloned placeholder is now the last shape in the slide
    return layout_placeholder, slide.shapes[-1]


def get_textframe(shape):
    """
    Get the shape's text frame and enable automatic resizing. The resize only
    triggers after opening the file in the PowerPoint application and making a change or saving.
    """
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return text_frame


def write_bullet(text_frame, text, level):
    """Write a bullet to the provided text frame at the specified level."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level


def write_objective_list(text_frame, objectives):
    """Write a list of objectives to the provided text frame."""
    for obj in objectives:
        status = obj["status"]
        if obj["complete"]:
            status = "Achieved"
        write_bullet(text_frame, f"{obj['objective']} – {status}", 1)


def prepare_for_pptx(value):
    """Strip HTML and clear 0x0D characters to prepare text for notes slides."""
    try:
        if value:
            return BeautifulSoup(value, "lxml").text.replace("\x0D", "")
        return "N/A"
    except Exception:
        logger.exception("Failed parsing this value for PPTX: %s", value)
        return ""


def delete_paragraph(par):
    """
    Delete the specified paragraph.

    **Parameter**

    ``par``
        Paragraph to delete from the document
    """
    p = par._p
    parent_element = p.getparent()
    if parent_element is not None:
        parent_element.remove(p)
    else:
        logger.warning("Could not delete paragraph in because it had no parent element")
